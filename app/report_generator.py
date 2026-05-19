"""Professional PowerPoint report generation."""

from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.chart_utils import render_bar_chart
from app.models import AnalysisReport

# Brand
NAVY = RGBColor(0x0F, 0x27, 0x44)
NAVY_MID = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0xE1, 0x1D, 0x48)
ACCENT_LIGHT = RGBColor(0xFB, 0x71, 0x85)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x64, 0x74, 0x8B)
TEXT = RGBColor(0x33, 0x41, 0x55)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
_page = 0


class PowerPointReportGenerator:
    def generate(self, report: AnalysisReport) -> bytes:
        global _page
        _page = 0
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        self._cover(prs, report)
        self._executive(prs, report)
        self._channel_table(prs, report)
        self._subscriber_chart(prs, report)
        self._rankings(prs, report)
        self._videos(prs, report)
        self._topics(prs, report)
        self._posting(prs, report)
        self._engagement(prs, report)
        self._gaps(prs, report)
        self._recommendations(prs, report)
        self._scorecard(prs, report)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _blank(self, prs):
        global _page
        _page += 1
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _bg(self, slide, color=LIGHT):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _footer(self, slide, report: AnalysisReport):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), SLIDE_W, Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        left = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(5), Inches(0.3))
        lp = left.text_frame.paragraphs[0]
        lp.text = "VidIntel · Video Competitor Intelligence"
        lp.font.size = Pt(8)
        lp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        right = slide.shapes.add_textbox(Inches(8.5), Inches(7.15), Inches(1.2), Inches(0.3))
        rp = right.text_frame.paragraphs[0]
        rp.text = str(_page)
        rp.font.size = Pt(9)
        rp.font.color.rgb = WHITE
        rp.alignment = PP_ALIGN.RIGHT

    def _header(self, slide, title: str, subtitle: str = ""):
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT
        stripe.line.fill.background()

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), 0, Inches(9.88), Inches(1.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9), Inches(0.55))
        tp = tb.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(26)
        tp.font.bold = True
        tp.font.color.rgb = WHITE

        if subtitle:
            sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(9), Inches(0.32))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(11)
            sp.font.color.rgb = ACCENT_LIGHT

    def _cover(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide, NAVY)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()

        deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-1), Inches(4), Inches(4))
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_MID
        deco.line.fill.background()

        brand = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(3), Inches(0.4))
        bp = brand.text_frame.paragraphs[0]
        bp.text = "VIDINTEL"
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = ACCENT_LIGHT

        title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.4))
        tp = title.text_frame.paragraphs[0]
        tp.text = "Video Competitor\nIntelligence Report"
        tp.font.size = Pt(40)
        tp.font.bold = True
        tp.font.color.rgb = WHITE

        names = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(1))
        np = names.text_frame.paragraphs[0]
        np.text = "  vs  ".join([report.company] + report.competitors)
        np.font.size = Pt(18)
        np.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

        date_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(4), Inches(0.4))
        dp = date_box.text_frame.paragraphs[0]
        dp.text = report.generated_at.strftime("%B %d, %Y")
        dp.font.size = Pt(13)
        dp.font.color.rgb = ACCENT_LIGHT

        self._footer(slide, report)

    def _executive(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Executive Summary", f"Market leader: {report.leader}")
        self._footer(slide, report)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(9.1), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(8.6), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = report.executive_summary
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.line_spacing = 1.3

        leader = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(3.85), Inches(9.1), Inches(2.5))
        leader.fill.solid()
        leader.fill.fore_color.rgb = NAVY
        leader.line.fill.background()

        lb = slide.shapes.add_textbox(Inches(0.75), Inches(4.05), Inches(8.5), Inches(2.1))
        ltf = lb.text_frame
        ltf.word_wrap = True
        p1 = ltf.paragraphs[0]
        p1.text = f"MARKET LEADER: {report.leader.upper()}"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_LIGHT
        p2 = ltf.add_paragraph()
        p2.text = report.leader_reason
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p2.space_before = Pt(10)

    def _channel_table(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Channel Overview", "Subscriber count · Content library · Publishing velocity")
        self._footer(slide, report)
        headers = ["Company", "Subscribers", "Videos", "Total Views", "Uploads/Mo"]
        rows = [[
            ins.company_name[:18], f"{ins.subscribers:,}", f"{ins.total_videos:,}",
            f"{ins.total_views:,}", f"{ins.uploads_per_month:.1f}",
        ] for ins in report.insights]
        self._table(slide, 0.45, 1.35, headers, rows)

    def _subscriber_chart(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Audience Reach", "Subscriber comparison across brands")
        self._footer(slide, report)
        m = report.comparative_metrics
        path = render_bar_chart(m["companies"], [float(s) for s in m["subscribers"]], "Subscribers by Brand")
        slide.shapes.add_picture(path, Inches(0.35), Inches(1.2), width=Inches(9.3))
        Path(path).unlink(missing_ok=True)

    def _rankings(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Competitive Rankings", "Composite video marketing score (0–100)")
        self._footer(slide, report)

        for i, r in enumerate(report.rankings[:5]):
            col = i % 3
            row = i // 3
            left = 0.45 + col * 3.15
            top = 1.35 + row * 2.85
            is_first = r.rank == 1

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.95), Inches(2.55))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB) if is_first else WHITE
            card.line.color.rgb = GOLD if is_first else RGBColor(0xE2, 0xE8, 0xF0)

            tb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.12), Inches(2.65), Inches(0.35))
            rp = tb.text_frame.paragraphs[0]
            rp.text = f"#{r.rank}  {r.company_name[:16]}"
            rp.font.size = Pt(13)
            rp.font.bold = True
            rp.font.color.rgb = GOLD if is_first else NAVY

            sc = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(2.65), Inches(0.7))
            sp = sc.text_frame.paragraphs[0]
            sp.text = f"{r.overall_score:.0f}"
            sp.font.size = Pt(36)
            sp.font.bold = True
            sp.font.color.rgb = ACCENT if is_first else NAVY

            # Score bar
            bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.15), Inches(top + 1.35), Inches(2.65), Inches(0.18))
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            bar_bg.line.fill.background()
            bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.15), Inches(top + 1.35), Inches(2.65 * r.overall_score / 100), Inches(0.18))
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = ACCENT if is_first else BLUE
            bar_fill.line.fill.background()

            metrics = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 1.65), Inches(2.65), Inches(0.8))
            mp = metrics.text_frame.paragraphs[0]
            mp.text = f"Aud {r.subscriber_score:.0f} · Eng {r.engagement_score:.0f} · Con {r.consistency_score:.0f}"
            mp.font.size = Pt(8)
            mp.font.color.rgb = SLATE

    def _videos(self, prs, report: AnalysisReport):
        for ins in report.insights[:5]:
            slide = self._blank(prs)
            self._bg(slide)
            self._header(slide, "Top Performing Videos", ins.company_name)
            self._footer(slide, report)

            for vi, v in enumerate(ins.top_videos[:3]):
                top = 1.35 + vi * 1.85
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(top), Inches(9.1), Inches(1.65))
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

                rank = slide.shapes.add_textbox(Inches(0.6), Inches(top + 0.15), Inches(0.5), Inches(0.4))
                rpp = rank.text_frame.paragraphs[0]
                rpp.text = str(vi + 1)
                rpp.font.size = Pt(22)
                rpp.font.bold = True
                rpp.font.color.rgb = ACCENT

                tb = slide.shapes.add_textbox(Inches(1.1), Inches(top + 0.15), Inches(7.5), Inches(0.9))
                tf = tb.text_frame
                tf.word_wrap = True
                tp = tf.paragraphs[0]
                title = v.title[:70] + ("…" if len(v.title) > 70 else "")
                tp.text = title
                tp.font.size = Pt(12)
                tp.font.bold = True
                tp.font.color.rgb = NAVY

                stats = slide.shapes.add_textbox(Inches(1.1), Inches(top + 0.95), Inches(8), Inches(0.45))
                sp = stats.text_frame.paragraphs[0]
                sp.text = f"{v.views:,} views  ·  {v.engagement_rate}% engagement  ·  {v.likes:,} likes  ·  {v.comments:,} comments"
                sp.font.size = Pt(10)
                sp.font.color.rgb = SLATE

    def _topics(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Content Topics & Themes", "Recurring themes from video titles")
        self._footer(slide, report)

        for idx, ins in enumerate(report.insights):
            col = idx % 2
            row = idx // 2
            left = 0.45 + col * 4.7
            top = 1.35 + row * 2.6
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.5), Inches(2.35))
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

            hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(4.5), Inches(0.45))
            hb.fill.solid()
            hb.fill.fore_color.rgb = NAVY
            hb.line.fill.background()

            ht = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.08), Inches(4.2), Inches(0.35))
            hp = ht.text_frame.paragraphs[0]
            hp.text = ins.company_name
            hp.font.size = Pt(12)
            hp.font.bold = True
            hp.font.color.rgb = WHITE

            topics = ins.top_topics[:6] or ["No recurring themes"]
            tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(4.1), Inches(1.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for ti, t in enumerate(topics):
                p = tf.paragraphs[0] if ti == 0 else tf.add_paragraph()
                p.text = f"• {t}"
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT
                p.space_after = Pt(4)

    def _posting(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Posting Frequency & Consistency", "Publishing cadence analysis")
        self._footer(slide, report)
        m = report.comparative_metrics
        path = render_bar_chart(m["companies"], m["uploads_per_month"], "Uploads per Month", horizontal=False)
        slide.shapes.add_picture(path, Inches(0.35), Inches(1.15), width=Inches(5.8))
        Path(path).unlink(missing_ok=True)

        for i, ins in enumerate(report.insights):
            top = 1.4 + i * 0.95
            row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(top), Inches(3.15), Inches(0.8))
            row.fill.solid()
            row.fill.fore_color.rgb = WHITE
            row.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            tb = slide.shapes.add_textbox(Inches(6.55), Inches(top + 0.12), Inches(2.9), Inches(0.6))
            tp = tb.text_frame.paragraphs[0]
            tp.text = f"{ins.company_name}\n{ins.uploads_per_month:.1f}/mo · Consistency {ins.posting_consistency_score:.0f}/100"
            tp.font.size = Pt(10)
            tp.font.color.rgb = TEXT

    def _engagement(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Engagement Analysis", "Average performance per video")
        self._footer(slide, report)
        m = report.comparative_metrics
        path = render_bar_chart(m["companies"], m["avg_engagement"], "Engagement Rate %", horizontal=True)
        slide.shapes.add_picture(path, Inches(0.35), Inches(1.15), width=Inches(5.5))
        Path(path).unlink(missing_ok=True)
        headers = ["Company", "Avg Views", "Likes", "Comments", "Eng %"]
        rows = [[
            ins.company_name[:14], f"{ins.avg_views_per_video:,.0f}",
            f"{ins.avg_likes_per_video:.0f}", f"{ins.avg_comments_per_video:.0f}",
            f"{ins.avg_engagement_rate:.2f}%",
        ] for ins in report.insights]
        self._table(slide, 6.1, 1.35, headers, rows, col_widths=[1.1, 0.85, 0.7, 0.75, 0.65])

    def _gaps(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Gap Analysis", f"Opportunities for {report.company}")
        self._footer(slide, report)
        gaps = report.gap_analysis[:5]
        if not gaps:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(1))
            tb.text_frame.paragraphs[0].text = "No significant content gaps detected."
            return
        for i, gap in enumerate(gaps):
            top = 1.3 + i * 1.1
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(top), Inches(9.1), Inches(0.95))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xEC, 0xFE, 0xFF)
            card.line.color.rgb = RGBColor(0x67, 0xE8, 0xF9)
            tb = slide.shapes.add_textbox(Inches(0.65), Inches(top + 0.1), Inches(8.7), Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = gap.topic_or_format
            p1.font.size = Pt(11)
            p1.font.bold = True
            p1.font.color.rgb = NAVY
            p2 = tf.add_paragraph()
            opp = gap.opportunity[:140] + ("…" if len(gap.opportunity) > 140 else "")
            p2.text = opp
            p2.font.size = Pt(9)
            p2.font.color.rgb = SLATE

    def _recommendations(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide)
        self._header(slide, "Strategic Recommendations", f"Action plan for {report.company}")
        self._footer(slide, report)
        for i, rec in enumerate(report.recommendations[:6]):
            top = 1.3 + i * 0.92
            num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(top), Inches(0.42), Inches(0.42))
            num.fill.solid()
            num.fill.fore_color.rgb = NAVY
            num.line.fill.background()
            nt = slide.shapes.add_textbox(Inches(0.45), Inches(top + 0.05), Inches(0.42), Inches(0.35))
            np = nt.text_frame.paragraphs[0]
            np.text = str(i + 1)
            np.font.size = Pt(12)
            np.font.bold = True
            np.font.color.rgb = WHITE
            np.alignment = PP_ALIGN.CENTER
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(8.5), Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = rec
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT

    def _scorecard(self, prs, report: AnalysisReport):
        slide = self._blank(prs)
        self._bg(slide, NAVY)
        self._header(slide, "Final Scorecard", "Overall video marketing ranking")
        self._footer(slide, report)
        headers = ["#", "Company", "Score", "Audience", "Engagement", "Consistency", "Volume"]
        rows = [[
            str(r.rank), r.company_name[:20], f"{r.overall_score:.0f}",
            f"{r.subscriber_score:.0f}", f"{r.engagement_score:.0f}",
            f"{r.consistency_score:.0f}", f"{r.content_volume_score:.0f}",
        ] for r in report.rankings]
        self._table(slide, 0.55, 1.45, headers, rows, on_dark=True)

    def _table(self, slide, left, top, headers, rows, col_widths=None, on_dark=False):
        n_cols = len(headers)
        n_rows = len(rows) + 1
        w = sum(col_widths) if col_widths else 9.0
        ts = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(w), Inches(0.38 * n_rows))
        table = ts.table
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = Inches(cw)
        for j, h in enumerate(headers):
            c = table.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = ACCENT if on_dark else NAVY
            p = c.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = table.cell(i + 1, j)
                c.text = val
                if i % 2 == 0:
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9) if not on_dark else NAVY_MID
                p = c.text_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT if not on_dark else WHITE
                p.alignment = PP_ALIGN.CENTER
